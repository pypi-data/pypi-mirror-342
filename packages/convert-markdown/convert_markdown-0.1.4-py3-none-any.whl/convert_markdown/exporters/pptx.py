from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
import base64
import io

class PPTXExporter:
    def __init__(self):
        self.processed_tables = set()
        
    def export(self, html: str) -> bytes:
        """Convert HTML to PPTX
        
        Args:
            html: Input HTML string
            
        Returns:
            bytes: PPTX presentation as bytes
        """
        try:
            soup = BeautifulSoup(html, "html.parser")
            prs = Presentation()
            
            # Set up layouts
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = soup.find('h1')
            if title:
                title_slide.shapes.title.text = title.get_text(strip=True)
            
            content_layout = prs.slide_layouts[1]  # Title and content
            blank_layout = prs.slide_layouts[6]    # Blank layout for images
            
            # Process content
            elements = soup.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'ol', 'table', 'img'])
            current_slide = None
            current_content = []
            # print("Processing elements: ", [element.name for element in elements])
            for element in elements:
                # Skip the title we already processed
                if element.name == 'h1' and element == title:
                    continue
                    
                # Start new slide for headings
                if element.name in ['h1', 'h2', 'h3']:
                    # Add previous slide content if exists
                    if current_content and current_slide:
                        self._add_content_to_slide(current_slide, current_content)
                    
                    current_slide = prs.slides.add_slide(content_layout)
                    current_slide.shapes.title.text = element.get_text(strip=True)
                    current_content = []
                
                # Handle different content types
                elif element.name in ['p', 'ul', 'ol']:
                    if not current_slide:
                        current_slide = prs.slides.add_slide(content_layout)
                    current_content.append(element.get_text(strip=True))
                
                elif element.name == 'table':
                    # print("Adding table slide")
                    if current_content and current_slide:
                        self._add_content_to_slide(current_slide, current_content)
                        current_content = []
                    
                    self._add_table_slide(prs, element)
                
                elif element.name == 'img':
                    if current_content and current_slide:
                        self._add_content_to_slide(current_slide, current_content)
                        current_content = []
                    
                    self._add_image_slide(prs, element)
            
            # Add any remaining content
            if current_content and current_slide:
                self._add_content_to_slide(current_slide, current_content)
            
            
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.read()
            
        except Exception as e:
            print(f"Error converting to PPTX: {e}")
            raise
            
    def _add_content_to_slide(self, slide, content):
        """Add text content to slide"""
        if not content:
            return
            
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = '\n\n'.join(content)
        
    def _add_table_slide(self, prs, table_elem):
        """Add table to new slide"""
        rows = table_elem.find_all('tr')
        if not rows:
            return
            
        # Create table hash to avoid duplicates
        table_hash = self._get_table_hash(rows)
        if table_hash in self.processed_tables:
            return
        self.processed_tables.add(table_hash)
        
        # Create slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        prev_heading = table_elem.find_previous(['h1', 'h2', 'h3'])
        slide.shapes.title.text = (prev_heading.get_text(strip=True) if prev_heading else "Table")
        
        # Add table
        cols = rows[0].find_all(['th', 'td'])
        shape_table = slide.shapes.add_table(
            rows=len(rows),
            cols=len(cols),
            left=Inches(1),
            top=Inches(2),
            width=Inches(8),
            height=Inches(0.5 * len(rows))
        ).table
        
        # Fill table
        for i, row in enumerate(rows):
            cells = row.find_all(['th', 'td'])
            for j, cell in enumerate(cells):
                shape_table.cell(i, j).text = cell.get_text(strip=True)
                if i == 0 or cell.name == 'th':
                    shape_table.cell(i, j).text_frame.paragraphs[0].font.bold = True
                    
    def _add_image_slide(self, prs, img_elem):
        """Add image to new slide"""
        if 'src' not in img_elem.attrs:
            return
            
        src = img_elem['src']
        if not src.startswith('data:image'):
            return
            
        # Create slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add image
        base64_str = src.split(',')[1]
        image_data = base64.b64decode(base64_str)
        image_stream = io.BytesIO(image_data)
        
        slide.shapes.add_picture(
            image_stream,
            left=Inches(1),
            top=Inches(1),
            width=Inches(8)
        )
        
        # Add title if it's a chart
        chart_container = img_elem.find_parent(class_='chart-container')
        if chart_container:
            prev_heading = chart_container.find_previous(['h1', 'h2', 'h3'])
            title_text = prev_heading.get_text(strip=True) if prev_heading else "Chart"
            
            txBox = slide.shapes.add_textbox(
                left=Inches(1),
                top=Inches(0.5),
                width=Inches(8),
                height=Inches(0.5)
            )
            txBox.text_frame.text = title_text
            
    def _get_table_hash(self, rows):
        """Create hash of table content to detect duplicates"""
        table_content = []
        for row in rows:
            cells = row.find_all(['th', 'td'])
            row_content = [cell.get_text(strip=True) for cell in cells]
            table_content.append('|'.join(row_content))
        return '||'.join(table_content)
        
    def _cleanup_empty_slides(self, prs):
        """Remove empty slides"""
        slides_to_remove = []
        for index, slide in enumerate(prs.slides):
            if index == 0:  # Keep title slide
                continue
                
            is_empty = True
            for shape in slide.shapes:
                if shape.shape_type == 15:  # Table
                    is_empty = False
                    break
                elif shape.shape_type == 13:  # Picture
                    is_empty = False
                    break
                elif shape.has_text_frame:
                    text = shape.text.strip()
                    if text and shape != slide.shapes.title:
                        is_empty = False
                        break
                        
            if is_empty:
                slides_to_remove.append(index)
                
        # Remove empty slides from end to start
        for index in reversed(slides_to_remove):
            xml_slides = prs.slides._sldIdLst
            xml_slides.remove(xml_slides[index]) 